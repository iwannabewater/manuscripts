#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Build the LLM RL algorithms slide deck.

The deck intentionally keeps formulas as first-class assets: every displayed
equation is compiled by LaTeX and converted to SVG via dvisvgm. The same formula
source is also converted to PNG for PPTX export.
"""

from __future__ import annotations

import html
import os
import re
import shutil
import subprocess
from datetime import date
from pathlib import Path
from typing import Any


ROOT = Path(__file__).resolve().parent
SLUG = "llm-rl-algorithms-2026"
TITLE = "大模型强化学习算法全景"
SUBTITLE = "从 RLHF、DPO 到 GRPO、DAPO、GSPO 的公式化理解与工程判据"
TODAY = "2026-05-16"


SOURCES: list[dict[str, str]] = [
    {"id": "TRPO", "title": "Trust Region Policy Optimization", "url": "https://arxiv.org/abs/1502.05477", "used": "TRPO 约束策略改进与 PPO 背景"},
    {"id": "GAE", "title": "High-Dimensional Continuous Control Using Generalized Advantage Estimation", "url": "https://arxiv.org/abs/1506.02438", "used": "GAE 优势估计公式"},
    {"id": "PPO", "title": "Proximal Policy Optimization Algorithms", "url": "https://arxiv.org/abs/1707.06347", "used": "PPO clipped surrogate"},
    {"id": "DRHF", "title": "Deep Reinforcement Learning from Human Preferences", "url": "https://arxiv.org/abs/1706.03741", "used": "早期人类偏好奖励建模"},
    {"id": "IGPT", "title": "Training language models to follow instructions with human feedback", "url": "https://arxiv.org/abs/2203.02155", "used": "InstructGPT 的 SFT、RM、PPO 流程"},
    {"id": "HH", "title": "Training a Helpful and Harmless Assistant with Reinforcement Learning from Human Feedback", "url": "https://arxiv.org/abs/2204.05862", "used": "Anthropic helpful harmless RLHF"},
    {"id": "CAI", "title": "Constitutional AI: Harmlessness from AI Feedback", "url": "https://arxiv.org/abs/2212.08073", "used": "RLAIF 与宪法反馈"},
    {"id": "RRHF", "title": "RRHF: Rank Responses to Align Language Models with Human Feedback", "url": "https://arxiv.org/abs/2304.05302", "used": "排序式 SFT 类方法"},
    {"id": "RAFT", "title": "RAFT: Reward rAnked FineTuning for Generative Foundation Model Alignment", "url": "https://arxiv.org/abs/2304.06767", "used": "best-of-K 奖励排序微调"},
    {"id": "SLIC", "title": "SLiC-HF: Sequence Likelihood Calibration with Human Feedback", "url": "https://arxiv.org/abs/2305.10425", "used": "margin ranking 与 CE 正则"},
    {"id": "DPO", "title": "Direct Preference Optimization", "url": "https://arxiv.org/abs/2305.18290", "used": "DPO 推导、隐式奖励、BT 损失"},
    {"id": "REST", "title": "Reinforced Self-Training", "url": "https://arxiv.org/abs/2308.08998", "used": "ReST 生成、过滤、微调循环"},
    {"id": "REMAX", "title": "ReMax: A Simple, Effective, and Efficient Reinforcement Learning Method for Aligning Large Language Models", "url": "https://arxiv.org/abs/2310.10505", "used": "贪心解码基线的 critic-free PG"},
    {"id": "IPO", "title": "A General Theoretical Paradigm to Understand Learning from Human Preferences", "url": "https://arxiv.org/abs/2310.12036", "used": "IPO 与 ΨPO 框架"},
    {"id": "KTO", "title": "KTO: Model Alignment as Prospect Theoretic Optimization", "url": "https://arxiv.org/abs/2402.01306", "used": "KTO 二元反馈与前景理论价值函数"},
    {"id": "GRPO", "title": "DeepSeekMath: Pushing the Limits of Mathematical Reasoning in Open Language Models", "url": "https://arxiv.org/abs/2402.03300", "used": "GRPO 定义与数学推理 RL"},
    {"id": "RLOO", "title": "Back to Basics: Revisiting REINFORCE Style Optimization for Learning from Human Feedback in LLMs", "url": "https://arxiv.org/abs/2402.14740", "used": "RLOO 与 REINFORCE-style RLHF"},
    {"id": "ORPO", "title": "ORPO: Monolithic Preference Optimization without Reference Model", "url": "https://arxiv.org/abs/2403.07691", "used": "odds-ratio preference objective"},
    {"id": "NASH", "title": "Nash Learning from Human Feedback", "url": "https://arxiv.org/abs/2403.08635", "used": "博弈与 mirror descent 视角"},
    {"id": "REWARDBENCH", "title": "RewardBench: Evaluating Reward Models for Language Modeling", "url": "https://arxiv.org/abs/2403.13787", "used": "奖励模型评估风险"},
    {"id": "SIMPO", "title": "SimPO: Simple Preference Optimization with a Reference-Free Reward", "url": "https://arxiv.org/abs/2405.14734", "used": "长度归一化 reference-free DPO 变体"},
    {"id": "ONLINE-DPO", "title": "Online DPO", "url": "https://arxiv.org/abs/2406.05534", "used": "在线采样偏好优化"},
    {"id": "TPO", "title": "Thought Preference Optimization", "url": "https://arxiv.org/abs/2410.10630", "used": "thought-level preference"},
    {"id": "TULU3", "title": "Tulu 3: Pushing Frontiers in Open Language Model Post-Training", "url": "https://arxiv.org/abs/2411.15124", "used": "开放后训练 recipe"},
    {"id": "RPP", "title": "REINFORCE++: Stabilizing Critic-Free Policy Optimization with Global Normalization", "url": "https://arxiv.org/abs/2501.03262", "used": "全局归一化 critic-free PPO"},
    {"id": "KIMI", "title": "Kimi k1.5: Scaling Reinforcement Learning with LLMs", "url": "https://arxiv.org/abs/2501.12599", "used": "长上下文、长思维链 RL recipe"},
    {"id": "R1", "title": "DeepSeek-R1: Incentivizing Reasoning Capability in LLMs via Reinforcement Learning", "url": "https://arxiv.org/abs/2501.12948", "used": "R1-Zero、rule-based rewards、两阶段 RL"},
    {"id": "DAPO", "title": "DAPO: An Open-Source LLM Reinforcement Learning System at Scale", "url": "https://arxiv.org/abs/2503.14476", "used": "Clip-Higher、Dynamic Sampling、token-level loss、overlong shaping"},
    {"id": "DRGRPO", "title": "Understanding R1-Zero-Like Training: A Critical Perspective", "url": "https://arxiv.org/abs/2503.20783", "used": "Dr.GRPO 对长度偏置与难度偏置的修正"},
    {"id": "VAPO", "title": "VAPO: Efficient and Reliable Reinforcement Learning for Advanced Reasoning Tasks", "url": "https://arxiv.org/abs/2504.05118", "used": "value-model augmented PPO 与长 CoT 稳定化"},
    {"id": "MINIMAX", "title": "MiniMax-M1: Scaling Test-Time Compute Efficiently with Lightning Attention", "url": "https://arxiv.org/abs/2506.13585", "used": "CISPO 与 off-policy 多轮更新"},
    {"id": "GSPO", "title": "Group Sequence Policy Optimization", "url": "https://arxiv.org/abs/2507.18071", "used": "sequence-level ratio 与 Qwen3 MoE 稳定训练"},
    {"id": "NCA", "title": "Noise Contrastive Alignment of Language Models with Explicit Rewards", "url": "https://arxiv.org/abs/2402.05369", "used": "NCA/InfoNCA 与显式标量奖励偏好优化"},
    {"id": "BCO", "title": "Binary Classifier Optimization for Large Language Model Alignment", "url": "https://arxiv.org/abs/2404.04656", "used": "二元 thumbs-up/down 反馈优化"},
    {"id": "SPIN", "title": "Self-Play Fine-Tuning Converts Weak Language Models to Strong Language Models", "url": "https://arxiv.org/abs/2401.01335", "used": "self-play fine-tuning 与自我改进"},
    {"id": "SPPO", "title": "Self-Play Preference Optimization for Language Model Alignment", "url": "https://arxiv.org/abs/2405.00675", "used": "自博弈偏好优化与 Nash 均衡"},
    {"id": "SAPO", "title": "Self-Augmented Preference Optimization: Off-Policy Paradigms for Language Model Alignment", "url": "https://arxiv.org/abs/2405.20830", "used": "off-policy self-augmented preference optimization"},
    {"id": "AOT", "title": "Distributional Preference Alignment of LLMs via Optimal Transport", "url": "https://arxiv.org/abs/2406.05882", "used": "optimal transport 与 distribution-level preference alignment"},
    {"id": "APO", "title": "Anchored Preference Optimization and Contrastive Revisions", "url": "https://arxiv.org/abs/2408.06266", "used": "anchored preference optimization 与 AI revisions"},
    {"id": "CALDPO", "title": "Cal-DPO: Calibrated Direct Preference Optimization for Language Model Alignment", "url": "https://arxiv.org/abs/2412.14516", "used": "implicit reward calibration"},
]


FORMULAS: dict[str, str] = {
    "mdp": r"\begin{aligned} V^\pi(s)&=\mathbb{E}_\pi\!\left[\sum_{t=0}^{\infty}\gamma^t r_t\mid s_0=s\right],\\ Q^\pi(s,a)&=\mathbb{E}_\pi\!\left[\sum_{t=0}^{\infty}\gamma^t r_t\mid s_0=s,a_0=a\right],\\ A^\pi(s,a)&=Q^\pi(s,a)-V^\pi(s) \end{aligned}",
    "kl_rlhf": r"\pi^\star=\arg\max_\pi\ \mathbb{E}_{x\sim\mathcal{D},\,y\sim\pi(\cdot|x)}[r_\phi(x,y)]-\beta\,D_{\mathrm{KL}}\!\left(\pi(\cdot|x)\,\|\,\pi_{\mathrm{ref}}(\cdot|x)\right)",
    "bt": r"\mathcal{L}_{\mathrm{RM}}(\phi)=-\mathbb{E}_{(x,y_w,y_l)}\log\sigma\!\left(r_\phi(x,y_w)-r_\phi(x,y_l)\right)",
    "pg": r"\nabla_\theta J(\theta)=\mathbb{E}_{\tau\sim\pi_\theta}\!\left[\sum_t\nabla_\theta\log\pi_\theta(a_t|s_t)\left(G_t-b(s_t)\right)\right]",
    "gae": r"\delta_t=r_t+\gamma V_\psi(s_{t+1})-V_\psi(s_t),\qquad \hat A_t^{\mathrm{GAE}(\gamma,\lambda)}=\sum_{l=0}^{\infty}(\gamma\lambda)^l\delta_{t+l}",
    "trpo": r"\max_\theta\ \mathbb{E}_t\!\left[\frac{\pi_\theta(a_t|s_t)}{\pi_{\theta_{\mathrm{old}}}(a_t|s_t)}\hat A_t\right]\quad\mathrm{s.t.}\quad \mathbb{E}_t[D_{\mathrm{KL}}(\pi_{\theta_{\mathrm{old}}}\|\pi_\theta)]\le \delta",
    "ppo": r"L_{\mathrm{CLIP}}(\theta)=\mathbb{E}_t\left[\min\left(r_t(\theta)\hat A_t,\ \mathrm{clip}(r_t(\theta),1-\epsilon,1+\epsilon)\hat A_t\right)\right]",
    "ppo_llm": r"\mathcal{L}_{\mathrm{PPO}}=\mathcal{L}_{\mathrm{clip}}-c_v\mathcal{L}_{V}+c_H\mathcal{H}(\pi_\theta)-\beta\,\widehat D_{\mathrm{KL}}(\pi_\theta\|\pi_{\mathrm{ref}})",
    "rloo": r"\hat g_{\mathrm{RLOO}}=\frac{1}{K}\sum_{i=1}^K\left(R_i-\frac{1}{K-1}\sum_{j\ne i}R_j\right)\nabla_\theta\log\pi_\theta(y_i|x)",
    "remax": r"\hat A(x,y)=R(x,y)-R(x,\hat y_{\mathrm{greedy}}),\qquad \hat g=\hat A(x,y)\nabla_\theta\log\pi_\theta(y|x)",
    "reinforcepp": r"\hat A_i' = R_i-\frac{1}{G}\sum_{j=1}^{G}R_j,\qquad \hat A_i=\frac{\hat A_i'-\mu_{\mathcal{B}}}{\sigma_{\mathcal{B}}+\varepsilon}",
    "raft": r"y^\star=\arg\max_{y_j\sim\pi_\theta(\cdot|x),\,j\le K} r(x,y_j),\qquad \theta\leftarrow\arg\max_\theta\log\pi_\theta(y^\star|x)",
    "slic": r"\mathcal{L}_{\mathrm{SLiC}}=\max\!\left(0,\delta-\log\pi_\theta(y_w|x)+\log\pi_\theta(y_l|x)\right)-\lambda\log\pi_\theta(y_{\mathrm{ref}}|x)",
    "dpo_reward": r"r_\theta(x,y)=\beta\log\frac{\pi_\theta(y|x)}{\pi_{\mathrm{ref}}(y|x)}+\beta\log Z(x)",
    "dpo": r"\mathcal{L}_{\mathrm{DPO}}=-\mathbb{E}\log\sigma\!\left(\beta\log\frac{\pi_\theta(y_w|x)}{\pi_{\mathrm{ref}}(y_w|x)}-\beta\log\frac{\pi_\theta(y_l|x)}{\pi_{\mathrm{ref}}(y_l|x)}\right)",
    "ipo": r"\mathcal{L}_{\mathrm{IPO}}=\mathbb{E}\left[\left(\log\frac{\pi_\theta(y_w|x)}{\pi_{\mathrm{ref}}(y_w|x)}-\log\frac{\pi_\theta(y_l|x)}{\pi_{\mathrm{ref}}(y_l|x)}-\frac{1}{2\beta}\right)^2\right]",
    "kto": r"\begin{aligned} r_\theta(x,y)&=\log\frac{\pi_\theta(y|x)}{\pi_{\mathrm{ref}}(y|x)},\\ v(x,y)&=\begin{cases}\lambda_D\sigma(\beta(r_\theta(x,y)-z_0)),&y\in\mathcal{D}\\ \lambda_U\sigma(\beta(z_0-r_\theta(x,y))),&y\in\mathcal{U}\end{cases} \end{aligned}",
    "orpo": r"\mathcal{L}_{\mathrm{ORPO}}=\mathcal{L}_{\mathrm{SFT}}-\lambda\,\mathbb{E}\log\sigma\!\left(\log\frac{\mathrm{odds}_\theta(y_w|x)}{\mathrm{odds}_\theta(y_l|x)}\right),\quad \mathrm{odds}_\theta(y|x)=\frac{P_\theta(y|x)}{1-P_\theta(y|x)}",
    "simpo": r"\mathcal{L}_{\mathrm{SimPO}}=-\mathbb{E}\log\sigma\!\left(\frac{\beta}{|y_w|}\log\pi_\theta(y_w|x)-\frac{\beta}{|y_l|}\log\pi_\theta(y_l|x)-\gamma\right)",
    "pl": r"P(y_1\succ\cdots\succ y_K|x)=\prod_{k=1}^{K}\frac{\exp r_\phi(x,y_k)}{\sum_{j=k}^{K}\exp r_\phi(x,y_j)}",
    "grpo": r"\mathcal{J}_{\mathrm{GRPO}}=\mathbb{E}\!\left[\frac{1}{G}\sum_{i=1}^{G}\frac{1}{|o_i|}\sum_t\min\!\left(r_{i,t}\hat A_i,\mathrm{clip}(r_{i,t},1-\epsilon,1+\epsilon)\hat A_i\right)-\beta D_{\mathrm{KL}}\right]",
    "grpo_adv": r"\hat A_i=\frac{R_i-\mathrm{mean}(\{R_j\}_{j=1}^{G})}{\mathrm{std}(\{R_j\}_{j=1}^{G})+\varepsilon}",
    "drgrpo": r"\tilde A_i=R_i-\frac{1}{G}\sum_{j=1}^{G}R_j,\qquad \mathcal{J}_{\mathrm{Dr.GRPO}}\propto\mathbb{E}\left[\sum_{i=1}^{G}\sum_t \rho_{i,t}\tilde A_i\right]",
    "dapo": r"\mathcal{J}_{\mathrm{DAPO}}=\mathbb{E}\!\left[\frac{1}{\sum_i|o_i|}\sum_{i=1}^{G}\sum_t\min\!\left(r_{i,t}\hat A_i,\mathrm{clip}(r_{i,t},1-\epsilon_{\mathrm{low}},1+\epsilon_{\mathrm{high}})\hat A_i\right)\right]",
    "gspo": r"\mathcal{J}_{\mathrm{GSPO}}=\mathbb{E}\!\left[\frac{1}{G}\sum_i\min\!\left(s_i(\theta)\hat A_i,\mathrm{clip}(s_i(\theta),1-\epsilon,1+\epsilon)\hat A_i\right)\right]",
    "gspo_ratio": r"s_i(\theta)=\left(\frac{\pi_\theta(y_i|x)}{\pi_{\theta_{\mathrm{old}}}(y_i|x)}\right)^{1/|y_i|}=\exp\!\left(\frac{1}{|y_i|}\sum_t\log\frac{\pi_\theta(y_{i,t}|x,y_{i,<t})}{\pi_{\theta_{\mathrm{old}}}(y_{i,t}|x,y_{i,<t})}\right)",
    "cispo": r"\mathcal{J}_{\mathrm{CISPO}}=\mathbb{E}\!\left[\frac{1}{\sum_i|o_i|}\sum_i\sum_t \mathrm{sg}\!\left(\hat r_{i,t}\right)\hat A_i\log\pi_\theta(o_{i,t}|q,o_{i,<t})\right],\quad \hat r_{i,t}=\mathrm{clip}(r_{i,t},1-\epsilon_{\mathrm{low}}^{IS},1+\epsilon_{\mathrm{high}}^{IS})",
    "vapo": r"\hat A_t=\sum_{l=0}^{T-t-1}(\gamma\lambda_{\mathrm{policy}})^l\delta_{t+l},\qquad \lambda_{\mathrm{policy}}=1-\frac{\alpha}{\ell}",
    "prm": r"R_{\mathrm{PRM}}(x,y_{1:T})=\sum_{k=1}^{K}r_\phi(x,y_{1:t_k}),\qquad R_{\mathrm{ORM}}(x,y)=\mathbb{1}[\mathrm{verify}(x,y)]",
}


def esc(value: Any) -> str:
    return html.escape(str(value), quote=True)


LATIN_TOKEN = re.compile(r"(?<![A-Za-z0-9_])([A-Za-z][A-Za-z0-9+./_-]*)(?![A-Za-z0-9_])")


def fmt(value: Any) -> str:
    raw = str(value)
    parts: list[str] = []
    last = 0
    for match in LATIN_TOKEN.finditer(raw):
        parts.append(esc(raw[last : match.start()]))
        parts.append(f"<span class='latin'>{esc(match.group(1))}</span>")
        last = match.end()
    parts.append(esc(raw[last:]))
    return "".join(parts)


def run(cmd: list[str], cwd: Path | None = None, quiet: bool = False) -> subprocess.CompletedProcess[str]:
    return subprocess.run(cmd, cwd=cwd, text=True, stdout=subprocess.PIPE if quiet else None, stderr=subprocess.PIPE if quiet else None, check=True)


def render_latex_formula(key: str, formula: str) -> None:
    out_dir = ROOT / "assets" / "formulas"
    tmp_dir = out_dir / "_tex"
    tmp_dir.mkdir(parents=True, exist_ok=True)
    tex = rf"""\documentclass{{article}}
\pagestyle{{empty}}
\usepackage{{amsmath,amssymb,mathtools,bm}}
\usepackage[dvipsnames]{{xcolor}}
\usepackage[active,tightpage]{{preview}}
\PreviewEnvironment{{equation*}}
\begin{{document}}
\color[rgb]{{0.078,0.078,0.075}}
\begin{{equation*}}
{formula}
\end{{equation*}}
\end{{document}}
"""
    tex_path = tmp_dir / f"{key}.tex"
    tex_path.write_text(tex, encoding="utf-8")
    svg_path = out_dir / f"{key}.svg"
    png_path = out_dir / f"{key}.png"
    try:
        run(["latex", "-interaction=nonstopmode", "-halt-on-error", tex_path.name], cwd=tmp_dir, quiet=True)
        run(["dvisvgm", "--no-fonts", "--exact", f"{key}.dvi", "-o", str(svg_path)], cwd=tmp_dir, quiet=True)
        run(["magick", "-density", "240", str(svg_path), "-background", "none", str(png_path)], quiet=True)
    except Exception as exc:
        fallback_svg = f"""<svg xmlns="http://www.w3.org/2000/svg" width="1400" height="140" viewBox="0 0 1400 140">
<rect width="1400" height="140" fill="none"/>
<text x="18" y="72" font-size="24" font-family="JetBrains Mono, monospace" fill="#141413">{esc(formula)}</text>
</svg>"""
        svg_path.write_text(fallback_svg, encoding="utf-8")
        try:
            from PIL import Image, ImageDraw

            img = Image.new("RGBA", (1600, 180), (255, 255, 255, 0))
            draw = ImageDraw.Draw(img)
            draw.text((20, 70), formula[:220], fill=(20, 20, 19, 255))
            img.save(png_path)
        except Exception:
            pass
        print(f"[formula fallback] {key}: {exc}")


def render_formulas() -> None:
    for key, formula in FORMULAS.items():
        render_latex_formula(key, formula)
    tmp_dir = ROOT / "assets" / "formulas" / "_tex"
    if tmp_dir.exists():
        shutil.rmtree(tmp_dir)


def slide(
    title: str,
    eyebrow: str,
    lead: str = "",
    points: list[str] | None = None,
    formula: str | None = None,
    formula_caption: str = "",
    columns: list[tuple[str, list[str]]] | None = None,
    table: list[list[str]] | None = None,
    callout: str = "",
    source: str = "",
    kind: str = "standard",
) -> dict[str, Any]:
    return {
        "title": title,
        "eyebrow": eyebrow,
        "lead": lead,
        "points": points or [],
        "formula": formula,
        "formula_caption": formula_caption,
        "columns": columns or [],
        "table": table or [],
        "callout": callout,
        "source": source,
        "kind": kind,
    }


SLIDES: list[dict[str, Any]] = [
    slide(
        "大模型强化学习算法全景",
        "LLM RL ALGORITHMS",
        "从 RLHF、RLAIF、RLVR 到 DPO、PPO、GRPO、DAPO、GSPO、CISPO、VAPO 的公式化理解与工程判据",
        ["范围截至 2026-05-16，重点覆盖国际论文、国内大厂 recipe、开源训练系统与推理模型后训练实践。", "所有核心目标函数均以 LaTeX 编译为矢量公式，便于 PDF 与 PPTX 审阅。"],
        callout="主线问题：给定一个预训练语言模型，如何用人类偏好、AI 反馈或可验证奖励，稳定地改变整段文本分布而不摧毁已有能力。",
        kind="cover",
    ),
    slide(
        "LLM 强化学习不是单一算法，而是三类目标函数的合流",
        "全景框架",
        "同一批算法可以按学习信号分为显式奖励、隐式偏好、可验证奖励；按优化方式分为 on-policy、offline、iterative SFT。",
        columns=[
            ("显式奖励 RL", ["先训练或定义奖励函数，再用 PPO、RLOO、GRPO、DAPO、GSPO 等策略梯度更新。", "优点是可在线采样、可探索；代价是 rollout、KL、方差和分布漂移都要管理。"]),
            ("隐式偏好优化", ["DPO、IPO、KTO、ORPO、SimPO 把偏好对直接变成分类或校准损失。", "优点是简单稳定；代价是受离线数据覆盖、reference、长度偏置和偏好噪声影响。"]),
            ("可验证推理 RL", ["RLVR 用规则、程序、单测、数学答案或 verifier 直接给 reward。", "它推动了 DeepSeek-R1、Kimi k1.5、DAPO、CISPO、GSPO、VAPO 等长 CoT 系列。"]),
        ],
        source="IGPT, DPO, GRPO, R1, DAPO, GSPO",
    ),
    slide(
        "从 MDP 到语言模型，状态、动作和回报都被重新解释",
        "基础定义",
        "在经典 MDP 中，策略逐步选择动作；在 LLM 中，一个 token 是动作，一个 prefix 是状态，一条 response 是轨迹。",
        points=[
            "状态：s_t=(x,y_{<t})，包含 prompt 与已生成 token。",
            "动作：a_t=y_t，词表规模通常达到数万到数十万。",
            "轨迹：整段输出 y_{1:T}，奖励可以只在末端给出，也可以在过程步骤给出。",
            "优势函数 A 衡量某动作相对当前状态平均行为的增益，是几乎所有策略梯度算法的核心低方差信号。",
        ],
        formula="mdp",
        formula_caption="值函数、动作值函数与优势函数的标准定义。",
    ),
    slide(
        "语言模型后训练的核心约束是奖励最大化与分布保持的张力",
        "KL 正则化 RLHF",
        "如果只最大化奖励，策略会钻奖励模型或 verifier 的漏洞；如果 KL 太强，模型几乎不学习。",
        points=[
            "π_ref 通常是 SFT 模型或上一阶段 policy，用来保留语言能力、风格和安全边界。",
            "β 是对齐中的温度和刹车：β 越大，偏离 reference 的成本越高。",
            "DPO 的推导、PPO 的 KL penalty、GRPO 的 reference KL 都可看作这个约束目标的不同求解方式。",
        ],
        formula="kl_rlhf",
        formula_caption="后训练中的 KL 正则化策略优化目标。",
        source="IGPT, DPO, RLOO",
    ),
    slide(
        "奖励信号的来源决定算法可承受的噪声和探索方式",
        "Reward Taxonomy",
        table=[
            ["奖励来源", "典型算法", "优势", "主要风险"],
            ["人类偏好 RM", "PPO、RLOO、DPO", "覆盖主观质量、安全、有用性", "标注噪声、RM 过优化、分布外外推"],
            ["AI feedback", "RLAIF、Constitutional AI", "规模化快，原则可控", "裁判模型偏见、同源模型自证"],
            ["规则或程序 verifier", "GRPO、DAPO、GSPO、CISPO", "稀疏但可信，适合数学代码", "奖励稀疏、格式投机、长答案偏置"],
            ["过程奖励 PRM", "PRM-PPO、VAPO", "能给中间步骤信用分配", "步骤标注昂贵，verifier 误差会传播"],
        ],
        source="DRHF, CAI, R1, VAPO, REWARDBENCH",
    ),
    slide(
        "Bradley-Terry 模型把相对偏好变成可优化的奖励差",
        "Reward Modeling",
        "人类经常不能稳定给出绝对分数，但能比较两个回答；BT 假设胜率由奖励差的 sigmoid 决定。",
        points=[
            "RM 学到的是相对排序，不是物理意义上的真实效用。",
            "训练 RM 后再做 PPO 是经典 RLHF；把 RM 消去并代入 KL 最优解，就得到 DPO 类损失。",
            "当偏好来自多个候选排序时，可用 Plackett-Luce 或 listwise ranking，RRHF 和 RAFT 属于这一路线的简化工程化形态。",
        ],
        formula="bt",
        formula_caption="pairwise reward model 的负对数似然。",
        source="DRHF, IGPT, DPO",
    ),
    slide(
        "RLHF pipeline 的每一步都在限制下一步的可学习空间",
        "SFT → RM → RL",
        "SFT 决定模型能否产生可接受候选；RM 决定奖励几何；RL 决定如何在不失稳的情况下移动分布。",
        columns=[
            ("SFT", ["把 foundation model 拉到 instruction-following 流形。", "数据质量决定后续 RL 是否在合理邻域内搜索。"]),
            ("RM 或 verifier", ["把偏好、规则、单测或 AI critique 映射为标量。", "奖励要有区分度，也要尽量难被利用。"]),
            ("Policy optimization", ["PPO 系列在线探索，DPO 系列离线校准，GRPO 系列用 group baseline 降低价值网络成本。", "最终目标是改 response 分布，不只是提高训练集 reward。"]),
        ],
        source="IGPT, HH, CAI, R1",
    ),
    slide(
        "PPO 被采用是因为它用剪切 surrogate 近似 trust region",
        "PPO",
        "LLM RLHF 早期沿用 PPO，是因为它比 TRPO 简单，又能避免单次更新把策略推离采样分布太远。",
        points=[
            "r_t 是新旧策略概率比；当 r_t 超出 1±ε，梯度被剪切。",
            "正优势 token 的概率不能被一次性推得太高，负优势 token 的概率也不能被一次性压得太低。",
            "LLM 中 PPO 通常还叠加 KL penalty、value loss、entropy bonus、reward whitening 和 rollout buffer。",
        ],
        formula="ppo",
        formula_caption="PPO clipped surrogate objective。",
        source="TRPO, PPO, IGPT",
    ),
    slide(
        "TRPO 的约束优化解释了 PPO 中 KL、clip 和 trust region 的来源",
        "TRPO → PPO",
        "TRPO 明确限制平均 KL；PPO 用剪切概率比实现一阶近似，牺牲精确 trust region 换取工程简洁。",
        points=[
            "TRPO 需要二阶近似和约束求解，难以直接套到超大 LLM 训练。",
            "PPO 的 clip 不等价于 KL 约束，因此 LLM 训练仍常监控 actual KL 并动态调 β。",
            "当 rollout 与训练步数过多时，r_t 会快速离开 1 附近，PPO 的有效梯度会塌缩。",
        ],
        formula="trpo",
        formula_caption="TRPO 的受约束 surrogate 目标。",
        source="TRPO, PPO",
    ),
    slide(
        "GAE 与 value model 给 PPO 低方差，也带来显存、算力和偏差",
        "Advantage Estimation",
        "PPO 的稳定很大程度来自 critic；但在 LLM 中，value head 或 value model 通常昂贵且难训。",
        points=[
            "λ 接近 1 时偏差低但方差高；λ 小时依赖 value bootstrap，偏差更大。",
            "长 CoT 中 reward 可能只在末端出现，critic 要学跨数千 token 的信用分配。",
            "GRPO、RLOO、ReMax、REINFORCE++ 的共同动机之一，就是避免训练一个可靠 critic。",
        ],
        formula="gae",
        formula_caption="GAE 的 TD residual 与指数加权优势。",
        source="GAE, RLOO, RPP",
    ),
    slide(
        "LLM-PPO 的真实目标通常是四项损失的折中",
        "PPO-RLHF Loss",
        "工业实现同时管理 clipped surrogate、奖励、价值、熵、KL、长度和预训练能力。",
        points=[
            "value loss 学习 response 或 token 级回报，决定 advantage 是否可信。",
            "entropy bonus 防止早期塌缩，但过强会牺牲遵循性。",
            "KL penalty 是防 reward hacking 的主要保护，常按 observed KL 自适应调整。",
            "InstructGPT 还使用预训练混合项缓解 alignment tax，这启发了后续的 SFT/RL 混合训练。",
        ],
        formula="ppo_llm",
        formula_caption="LLM-RLHF 中常见的 PPO 合成目标示意。",
        source="IGPT, PPO",
    ),
    slide(
        "策略梯度本身并不复杂，难点在于基线、归一化和采样分布",
        "REINFORCE",
        "REINFORCE 是无 critic 的最基本策略梯度；在 SFT 起点足够强时，LLM 的方差问题比传统随机初始化 RL 更可控。",
        points=[
            "b(s_t) 不改变无偏性，但能显著降低方差。",
            "如果 reward 是 response-level，所有 token 共享同一回报，长度和格式会影响梯度规模。",
            "critic-free 方法的核心差异就是如何构造 baseline 与 normalization。",
        ],
        formula="pg",
        formula_caption="带 baseline 的 score-function 梯度。",
    ),
    slide(
        "RLOO 用同一 prompt 的其他样本作为 leave-one-out 基线",
        "RLOO",
        "RLOO 避免 value model，又比单样本 REINFORCE 方差更低，适合每个 prompt 采 K 个回答的 RLHF/RLVR 设置。",
        points=[
            "每个样本的 baseline 是同组其他样本平均 reward，因此不会用自己的 reward 抵消自己。",
            "RLOO 与 GRPO 的差别在于是否使用标准差归一化、PPO clip、KL 处理和 token-level objective。",
            "Back to Basics 证明在不少 RLHF 场景下，REINFORCE-style 方法可以达到或超过 PPO。",
        ],
        formula="rloo",
        formula_caption="leave-one-out baseline 的策略梯度估计。",
        source="RLOO",
    ),
    slide(
        "ReMax 用 greedy response 的奖励作为确定性基线",
        "ReMax",
        "ReMax 的思想是用当前策略的贪心输出估计 baseline，再对采样输出的相对改进做策略梯度。",
        points=[
            "每个 prompt 只需一个采样 response 和一个 greedy response，开销低。",
            "baseline 与 prompt 强相关，能减少不同 prompt 难度带来的 reward 尺度差。",
            "缺点是 greedy baseline 可能受解码策略影响，并且对 group 内多样性利用不如 RLOO/GRPO。",
        ],
        formula="remax",
        formula_caption="ReMax 的 greedy-baseline advantage。",
        source="REMAX",
    ),
    slide(
        "REINFORCE++ 把 critic-free 训练的关键放到全局归一化",
        "REINFORCE++",
        "它批评 GRPO/RLOO 的 prompt-local normalization 会引入难度偏置，并用 batch-level normalization 稳定优势尺度。",
        points=[
            "先用组均值减去 prompt 难度，再在全局 batch 上标准化。",
            "全局标准化让不同 prompt 的可学习信号在一个尺度上竞争。",
            "实现上仍可结合 PPO-style clip 与 reference KL loss，保持训练稳定。",
        ],
        formula="reinforcepp",
        formula_caption="REINFORCE++ 的组内去均值与全局标准化。",
        source="RPP",
    ),
    slide(
        "Best-of-N 和 rejection sampling 是最朴素的策略改进算子",
        "Sampling as Policy Improvement",
        "不更新参数也能通过采样与筛选提升输出质量；更新参数时，RAFT/ReST/Rejection SFT 把筛选结果蒸馏回模型。",
        points=[
            "Best-of-N 依赖 reward 排序，推理成本随 N 线性增长。",
            "Rejection sampling 把高 reward 样本收集为新 SFT 数据，稳定但探索慢。",
            "这类方法可看作硬 EM：E 步生成并选择候选，M 步最大似然拟合被选候选。",
        ],
        formula="raft",
        formula_caption="RAFT 的 reward-ranked fine-tuning 更新。",
        source="RAFT, REST, R1",
    ),
    slide(
        "DPO 的关键洞察是把 KL 最优策略反解成隐式奖励",
        "DPO Derivation",
        "在 KL 正则化 RLHF 中，最优策略与 reference 的 log ratio 等于奖励的仿射变换；于是偏好损失可以直接写在 policy 上。",
        points=[
            "不需要显式训练 reward model，不需要 rollout，不需要 value model。",
            "β 控制 policy 相对 reference 的更新强度，也决定隐式 reward 的尺度。",
            "DPO 本质是对 winner-loser log odds margin 的 logistic 校准。",
        ],
        formula="dpo_reward",
        formula_caption="DPO 的隐式奖励重参数化。",
        source="DPO",
    ),
    slide(
        "DPO 把偏好对学习写成 reference-corrected logistic loss",
        "DPO Objective",
        "赢家相对 reference 的 log probability 应该比输家提高；输家相对 reference 的 log probability 应该降低。",
        points=[
            "reference correction 避免把原模型已经偏好的回答误判为训练信号。",
            "DPO 的训练稳定来自 supervised-like loss，但它牺牲了在线探索和 reward shaping 的灵活性。",
            "当偏好数据可分且 β 较大时，DPO 可能过度推向确定性 winner。",
        ],
        formula="dpo",
        formula_caption="Direct Preference Optimization 的标准 pairwise loss。",
        source="DPO, IPO",
    ),
    slide(
        "DPO 的风险集中在分布外偏好、长度效应和隐式 reward 饱和",
        "DPO Failure Modes",
        "离线偏好优化只有训练对中的局部比较；它不保证新策略在自身采样分布下仍被偏好。",
        points=[
            "长度偏置：log probability 随 token 数累积，若不处理会偏好短答或模板化答案。",
            "分布漂移：训练数据来自旧策略，新策略的错误区域没有被偏好对覆盖。",
            "饱和：sigmoid loss 在大 margin 区域梯度小，错误校准可能被固定。",
            "多轮迭代 DPO、online DPO、SimPO、IPO、KTO 都在修补这些问题的不同侧面。",
        ],
        source="DPO, IPO, SIMPO, ONLINE-DPO",
    ),
    slide(
        "SLiC-HF 用 max-margin 约束校准序列似然",
        "SLiC-HF",
        "SLiC-HF 是 DPO 之前的重要桥梁：它直接要求 preferred response 的 sequence likelihood 高于 rejected response 一个 margin。",
        points=[
            "ranking loss 负责对齐偏好，CE 正则负责保留 reference answer 的语言建模能力。",
            "它无需显式 reward model，也不需要在线 RL。",
            "相对 DPO，SLiC-HF 没有从 KL-RLHF 最优解中得到同样完整的概率解释。",
        ],
        formula="slic",
        formula_caption="SLiC-HF 的 margin ranking 加 CE 正则。",
        source="SLIC",
    ),
    slide(
        "RRHF 把多个候选的奖励排序蒸馏成 log-prob 排序",
        "RRHF",
        "RRHF 让模型对高奖励回答赋予更高的平均 token log probability，无需估计绝对奖励。",
        points=[
            "数据结构通常是一题多答，每个回答有 reward 或排序。",
            "训练信号是 pairwise/listwise ranking violation，辅以 SFT 保持流畅性。",
            "它适合已有高质量候选集的场景，但缺乏在线探索机制。",
        ],
        formula="pl",
        formula_caption="多候选排序可用 Plackett-Luce 似然表达。",
        source="RRHF",
    ),
    slide(
        "IPO 用平方校准避免 DPO 在可分数据上无限推大 margin",
        "IPO",
        "IPO 属于更一般的 ΨPO 框架，用 identity link 代替 Bradley-Terry sigmoid，优化有限目标 margin。",
        points=[
            "DPO 间接假设 BT 偏好模型；IPO 认为偏好概率不一定要被 BT 完全解释。",
            "平方损失让 log-ratio margin 靠近 1/(2β)，而不是越大越好。",
            "当偏好数据噪声较大或可分性很强时，IPO 更像校准器而不是无界分类器。",
        ],
        formula="ipo",
        formula_caption="常见参数化下的 IPO loss。",
        source="IPO",
    ),
    slide(
        "KTO 把偏好学习从成对比较放宽到单样本好坏标签",
        "KTO",
        "KTO 只需要 desirable/undesirable 标签，适合日志、审核、弱反馈和非成对偏好数据。",
        points=[
            "价值函数来自 prospect theory，对收益和损失使用不同权重。",
            "z0 是 reference point，通常与 KL 或批内基准有关。",
            "λ_D 与 λ_U 可处理正负样本不平衡，这是 DPO 不天然具备的优势。",
        ],
        formula="kto",
        formula_caption="KTO 的隐式 reward 与价值函数。",
        source="KTO",
    ),
    slide(
        "ORPO 把 SFT 和偏好对齐合并为单阶段 reference-free objective",
        "ORPO",
        "ORPO 不使用 reference model，而是在 SFT likelihood 上追加 winner-loser odds ratio penalty。",
        points=[
            "正样本仍通过 NLL 学习，负样本通过 odds ratio 被显式压低。",
            "reference-free 简化训练内存，但失去 reference correction 的保护。",
            "适合希望从偏好数据直接做 instruction tuning 的场景。",
        ],
        formula="orpo",
        formula_caption="ORPO 的 SFT 与 odds-ratio penalty。",
        source="ORPO",
    ),
    slide(
        "SimPO 用长度归一化 log-prob 构造 reference-free reward",
        "SimPO",
        "SimPO 认为生成时实际比较的是平均 token log probability，因此把 reward 直接定义为长度归一化 log probability。",
        points=[
            "去掉 reference model，降低显存与实现复杂度。",
            "γ margin 防止 winner 只是略微好一点就被过度强化。",
            "长度归一化直接处理 DPO 类 loss 的长短回答尺度问题。",
        ],
        formula="simpo",
        formula_caption="SimPO 的 reference-free margin objective。",
        source="SIMPO",
    ),
    slide(
        "Online DPO 试图修复离线偏好优化的分布漂移",
        "Online Preference Optimization",
        "在线变体用当前 policy 采样候选，再由人类、AI judge、reward model 或 verifier 产生偏好，随后更新 policy。",
        points=[
            "优点是训练对来自当前策略，能持续覆盖新错误区域。",
            "代价是需要在线评审管线，且 judge 的偏差会被循环放大。",
            "本质上，它把 DPO 从静态 supervised loss 变成 rollout-and-rank 策略迭代。",
        ],
        source="ONLINE-DPO, REST, R1",
    ),
    slide(
        "APO、AOT、NCA、SPPO 等变体主要在偏好几何上重写 DPO",
        "Post-DPO Variants",
        table=[
            ["算法族", "核心变化", "适用直觉"],
            ["APO", "给 preference margin 加 anchor 或正则，防止 winner 与 loser 同向漂移", "需要更明确控制参考点的离线对齐"],
            ["AOT", "用分布匹配或 optimal transport 视角处理偏好排序", "希望优化整体分布而非独立 pair"],
            ["NCA/BCO", "把 winner-loser 区分改写为 contrastive 或 classifier objective", "偏好数据噪声大、需要分类校准"],
            ["SPPO/Nash", "用自博弈、Nash 或 mirror descent 解释 policy improvement", "多模型互评、AI feedback 可持续生成"],
        ],
        callout="这些方法共同改变 log-ratio margin 的几何、参考点或采样过程，并不要求全面替代 DPO。",
        source="NASH, APO, AOT, NCA, BCO, SPPO",
    ),
    slide(
        "DPO 变体索引显示，偏好优化正在从 pairwise margin 走向校准与分布对齐",
        "Preference Optimization Index",
        table=[
            ["算法", "训练信号", "相对 DPO 的变化"],
            ["Cal-DPO", "偏好对与隐式 reward 校准项", "不仅比较 winner 与 loser，还约束隐式 reward 的绝对尺度"],
            ["NCA/InfoNCA", "显式标量 reward 或偏好数据", "用 noise-contrastive estimation 统一 reward data 与 pairwise preference"],
            ["BCO", "单样本 thumbs-up/down 二元反馈", "把 alignment 写成 binary classifier，logit 可解释为 reward"],
            ["AOT", "正负样本分布或未配对偏好数据", "用 optimal transport 约束正样本 reward 分布一阶随机占优"],
            ["APO", "AI revisions 产生的 contrastive pair", "用 anchor 控制 winner 与 loser 的漂移方向"],
        ],
        source="CALDPO, NCA, BCO, AOT, APO",
    ),
    slide(
        "Self-play 系列把偏好优化解释为模型与自身或对手策略的博弈",
        "Self-Play Alignment",
        table=[
            ["算法", "核心循环", "工程含义"],
            ["SPIN", "当前模型生成负例，SFT 模型或数据作为正例", "不依赖新增人工偏好，强调自我改进"],
            ["SPPO", "用 preference model 驱动 self-play policy update", "目标是逼近 Nash equilibrium，处理偏好非传递性"],
            ["SAPO", "EMA policy 与 replay buffer 产生 off-policy 自增强数据", "提高数据探索与复用效率"],
            ["Nash-MD", "mirror descent 或博弈视角优化 preference objective", "把 RLHF 看成寻找稳定策略而非单向爬 reward"],
            ["ReST", "生成、筛选、微调反复迭代", "最稳定的 self-training baseline，常与 DPO/RLVR 混合"],
        ],
        source="SPIN, SPPO, SAPO, NASH, REST",
    ),
    slide(
        "ReST 与 SPIN 类方法把策略改进离散化为生成、筛选、微调循环",
        "Iterative Self-Training",
        "这条路线少用显式 RL optimizer，多用高质量候选蒸馏；它稳定、便宜，但依赖筛选器质量。",
        points=[
            "Grow：用当前模型为未标注 prompt 生成多个候选。",
            "Improve：按 reward、规则、judge 或 verifier 过滤高分样本。",
            "Update：用 SFT、DPO 或 ranking loss 微调模型，然后重复。",
            "DeepSeek-R1 的 rejection sampling 与蒸馏阶段，本质上也使用了这类思想。",
        ],
        source="REST, R1",
    ),
    slide(
        "RLVR 将奖励模型换成可验证规则，重新打开 on-policy RL 的价值",
        "RL with Verifiable Rewards",
        "数学、代码、逻辑题的答案可由规则、解析器、单测或符号验证器判断，减少 RM 外推风险。",
        points=[
            "reward 通常稀疏：对就是 1，错就是 0，过程信息很少。",
            "group sampling 可把同一题的多个解法变成相对优势。",
            "长 CoT 会自然出现探索、反思与自我纠错，但也会产生 overthinking 和长度偏置。",
        ],
        formula="prm",
        formula_caption="ORM 与 PRM 的奖励表达。",
        source="GRPO, R1, VAPO",
    ),
    slide(
        "GRPO 用组内相对优势去掉 value model",
        "GRPO",
        "DeepSeekMath 提出的 GRPO 针对每个 prompt 采样 G 个回答，用组均值和组标准差归一化 reward。",
        points=[
            "无需训练 critic，节省显存与训练不稳定源。",
            "组内标准化让不同难度题目的 reward 尺度更接近。",
            "PPO-style clip 仍限制每个 token 的更新幅度；reference KL 约束策略漂移。",
            "它特别适合答案可验证、每题可多采样的数学和代码任务。",
        ],
        formula="grpo",
        formula_caption="GRPO 的 PPO-style group objective。",
        source="GRPO, R1",
    ),
    slide(
        "GRPO 的优势估计来自组内标准化 reward",
        "GRPO Advantage",
        "同一 prompt 的多个回答共享题目难度，因此组内相对 reward 可以作为低成本 advantage。",
        points=[
            "如果一组全对或全错，标准差接近 0，训练信号会消失或不稳定。",
            "如果题目太容易，所有回答高分；如果太难，所有回答低分，二者都不贡献有效梯度。",
            "DAPO 的 Dynamic Sampling 正是显式过滤这类无效组。",
        ],
        formula="grpo_adv",
        formula_caption="GRPO 的 group-relative advantage。",
        source="GRPO, DAPO",
    ),
    slide(
        "DeepSeek-R1 证明 rule-based reward 能诱导长链推理行为",
        "DeepSeek-R1",
        "R1-Zero 从 base model 直接做 RL，使用准确性奖励与格式奖励；R1 则加入冷启动 SFT、两阶段 RL 与蒸馏。",
        points=[
            "准确性奖励来自数学答案、代码单测或规则验证，避免训练神经 RM 产生 reward hacking。",
            "格式奖励约束思维链与答案边界，让 verifier 更稳定。",
            "训练中出现自我反思和长 CoT，但也伴随语言混杂、可读性和长度控制问题。",
            "R1 的工程 recipe 说明纯 RL 可激发推理，SFT 与蒸馏仍对可用性关键。",
        ],
        source="R1",
    ),
    slide(
        "Dr.GRPO 指出 GRPO 的标准差和长度归一化会引入偏置",
        "Dr.GRPO",
        "Dr.GRPO 的批评集中在两点：按 response length 除会改变长短回答的梯度权重，组内 std 会把题目难度混入 advantage 尺度。",
        points=[
            "长度归一化可能鼓励模型通过改变输出长度而非提高正确率来获得更大更新。",
            "std 归一化让高难题与低难题的相对信号被重新缩放，引入 difficulty bias。",
            "Dr.GRPO 建议移除这些项，使目标更接近 unbiased Monte Carlo policy gradient。",
        ],
        formula="drgrpo",
        formula_caption="Dr.GRPO 的去标准差、去长度偏置形式。",
        source="DRGRPO",
    ),
    slide(
        "DAPO 是对 naive GRPO 的四项工程级修复",
        "DAPO",
        "DAPO 的名称对应 Decoupled Clip and Dynamic sAmpling Policy Optimization，目标是让开源大规模 RLVR 训练可复现、可扩展。",
        points=[
            "Clip-Higher：分离 ε_low 与 ε_high，给正优势 token 更大的上行空间。",
            "Dynamic Sampling：过滤全对或全错组，只训练有区分度的 prompt。",
            "Token-level loss：按全 batch token 总数归一化，避免每条样本等权带来的长度偏置。",
            "Overlong shaping：对超过长度预算的回答做平滑惩罚或过滤，缓解无效长 CoT。",
        ],
        formula="dapo",
        formula_caption="DAPO 的 decoupled clip 与 token-level objective。",
        source="DAPO",
    ),
    slide(
        "DAPO 通过提高有效梯度密度改进 GRPO",
        "DAPO Mechanics",
        "在大规模 RLVR 中，常见失败来自有效 batch 太少、熵过早塌缩和长样本吞噬 token budget。",
        points=[
            "Dynamic Sampling 提高每个 optimizer step 的非零 advantage 比例。",
            "Clip-Higher 让低概率但正确的 token 有机会被充分强化。",
            "Token-level normalization 让长回答按 token 贡献梯度，而不是每条回答天然等权。",
            "Overlong reward shaping 把长度从副作用变成可控约束。",
        ],
        source="DAPO",
    ),
    slide(
        "GSPO 把重要性采样比从 token level 提升到 sequence level",
        "GSPO",
        "GSPO 认为 GRPO 的 token-level ratio 与 response-level reward 不匹配，尤其会破坏 MoE 模型的稳定训练。",
        points=[
            "reward 是整段 response 的属性，优势也是 response-level；因此 clip 也应作用在 sequence ratio 上。",
            "长度归一化的 sequence ratio 避免长回答概率乘积过小。",
            "在 MoE 中，sequence-level 更新减少 token 级路由噪声带来的不稳定。",
        ],
        formula="gspo",
        formula_caption="GSPO 的 sequence-level clipped objective。",
        source="GSPO",
    ),
    slide(
        "GSPO 的 ratio 是整段 log-ratio 的长度归一化指数",
        "GSPO Ratio",
        "这使得一个回答内的 token 不再各自独立触发 clipping，而是作为一个整体被接受或限制。",
        points=[
            "当整段回答整体优于旧策略时，所有 token 共享同一 sequence-level 权重。",
            "这与 response-level reward 的语义一致，减少 token-level credit 噪声。",
            "如果需要 token-level 分析，可以把 sequence ratio 看成平均 log-ratio 的 exponentiated form。",
        ],
        formula="gspo_ratio",
        formula_caption="GSPO 的长度归一化 sequence importance ratio。",
        source="GSPO",
    ),
    slide(
        "CISPO 剪切 IS 权重而不是剪切 token 更新",
        "CISPO",
        "MiniMax-M1 在多轮 off-policy 更新中发现，PPO/GRPO clip 会丢掉许多稀有但关键的反思 token 梯度。",
        points=[
            "CISPO 使用 clipped IS weight 修正分布差异，但所有 token 仍保留 log-prob 梯度。",
            "它特别关注 Wait、However、Recheck、Aha 等低概率推理转折 token。",
            "适合每个 generation batch 复用多次的训练设置，但必须监控 off-policy drift。",
        ],
        formula="cispo",
        formula_caption="CISPO 的 clipped IS-weight policy objective。",
        source="MINIMAX",
    ),
    slide(
        "VAPO 说明 value-based PPO 仍能在长 CoT 上胜出",
        "VAPO",
        "VAPO 不是放弃 critic，而是把 critic 训练、GAE、token-level loss 和正样本 NLL 做成更可靠的组合。",
        points=[
            "Value warmup 与 value pretraining 缓解初期 advantage 噪声。",
            "Decoupled GAE 用不同 λ 服务 critic 与 policy，长序列下减少偏差和方差冲突。",
            "Length-adaptive GAE 让更长回答使用更接近 1 的 λ，保留长程 credit。",
            "Positive NLL loss 把正确样本继续作为 supervised anchor，降低 RL 破坏可读性的风险。",
        ],
        formula="vapo",
        formula_caption="VAPO 中 length-adaptive policy GAE 的核心形式。",
        source="VAPO",
    ),
    slide(
        "PRM 把奖励从最终答案推进到推理过程，但也把标注误差推进到每一步",
        "Process Reward Models",
        "ORM 只看最终正确性；PRM 评价中间步骤，理论上能解决长 CoT 信用分配，实际依赖步骤粒度和 verifier 可靠性。",
        points=[
            "数学推理中，PRM 可奖励关键变形、约束检查和中间结论。",
            "代码任务中，过程奖励可来自编译、单测、静态分析、覆盖率或 partial credit。",
            "PRM 的危险是局部正确不保证全局正确，模型可能学习迎合过程裁判的表面模式。",
        ],
        formula="prm",
        formula_caption="过程奖励与结果奖励的抽象表达。",
        source="VAPO, REWARDBENCH",
    ),
    slide(
        "RLAIF 与 Constitutional AI 把人类原则转化为可规模化裁判",
        "RLAIF",
        "Anthropic 的 Constitutional AI 显示，原则、critique、revision 和 AI feedback 可以替代一部分人类 harmlessness 标注。",
        points=[
            "AI critique 先让模型按原则修改答案，再用 AI preference 训练 harmlessness。",
            "优势是规模、速度和一致性；风险是裁判模型继承同源偏见。",
            "RLAIF 不改变优化器本质，仍可接 PPO、DPO、KTO 或 rejection SFT。",
        ],
        source="CAI, HH",
    ),
    slide(
        "算法选择首先看奖励是否可在线生成，其次看是否需要探索",
        "Algorithm Matrix",
        table=[
            ["家族", "代表", "需要 rollout", "需要 RM/critic", "最适合"],
            ["PPO/value-based", "PPO、VAPO", "需要", "通常需要 critic", "复杂奖励、需要稳定 trust region"],
            ["critic-free PG", "RLOO、ReMax、REINFORCE++", "需要", "不需要 critic", "可多采样、reward 噪声可控"],
            ["offline preference", "DPO、IPO、KTO、ORPO、SimPO", "不需要", "不需要 critic", "已有偏好数据、快速迭代"],
            ["group RLVR", "GRPO、DAPO、GSPO、CISPO", "需要", "不需要 critic", "数学代码、规则 reward、长 CoT"],
            ["iterative SFT", "RAFT、ReST、rejection SFT", "生成即可", "需要筛选器", "稳定蒸馏、高质量候选充足"],
        ],
    ),
    slide(
        "离线偏好优化的差异主要是 reference、长度归一化和目标 margin",
        "Offline Preference Matrix",
        table=[
            ["算法", "核心 loss", "reference", "关键修补点"],
            ["DPO", "BT logistic on log-ratio margin", "需要", "把 KL-RLHF 反解成 supervised loss"],
            ["IPO", "square margin calibration", "需要", "防止可分数据上无界推大 margin"],
            ["KTO", "prospect-theory binary loss", "需要", "不要求成对偏好，处理正负不平衡"],
            ["ORPO", "SFT + odds-ratio penalty", "不需要", "单阶段训练，显存更低"],
            ["SimPO", "length-normalized log-prob margin", "不需要", "对齐生成度量并显式处理长度"],
        ],
        source="DPO, IPO, KTO, ORPO, SIMPO",
    ),
    slide(
        "on-policy RL 家族的差异主要是 advantage、clip 和归一化",
        "On-Policy Matrix",
        table=[
            ["算法", "baseline/advantage", "clip 对象", "主要贡献"],
            ["PPO", "critic + GAE", "token ratio", "稳定通用 RLHF"],
            ["RLOO", "同组 leave-one-out", "可配 PPO clip", "去 critic 降方差"],
            ["GRPO", "组均值与组 std", "token ratio", "RLVR 中去 value model"],
            ["DAPO", "GRPO advantage", "低高 ε 解耦 token ratio", "提高有效梯度和长 CoT 可控性"],
            ["GSPO", "组相对 response advantage", "sequence ratio", "对齐 response reward 与 MoE 稳定"],
            ["CISPO", "group advantage", "IS weight", "多轮 off-policy 更新保留关键 token 梯度"],
        ],
        source="PPO, RLOO, GRPO, DAPO, GSPO, MINIMAX",
    ),
    slide(
        "OpenAI InstructGPT 的贡献是把对齐工程流程标准化",
        "Industry Recipe: OpenAI",
        "InstructGPT 不是提出新 optimizer，而是证明 SFT、偏好 RM、PPO、KL 与人工评测能把 GPT-3 系列变成可用助手。",
        points=[
            "SFT 数据提供 instruction-following 初始分布。",
            "labeler ranking 训练 reward model，PPO 在 RM 上优化。",
            "KL 与预训练混合减少 reward overoptimization 和 alignment tax。",
            "后续 ChatGPT、RLHF 标准范式和开源实现都沿用这一流程骨架。",
        ],
        source="IGPT",
    ),
    slide(
        "Anthropic 把 helpfulness 与 harmlessness 分开建模，并引入原则驱动反馈",
        "Industry Recipe: Anthropic",
        "Helpful-Harmless RLHF 与 Constitutional AI 共同强调：奖励不是单一质量分，而是价值约束的组合。",
        points=[
            "HH 数据用成对比较表达 helpful 与 harmless 目标。",
            "Constitutional AI 用原则产生 critique 与 revision，再用 AI feedback 训练。",
            "这条路线解释了为什么大模型安全对齐常需要多奖励、多阶段、多裁判。",
        ],
        source="HH, CAI",
    ),
    slide(
        "DeepSeek 的路线把 GRPO 与可验证奖励推到推理模型中心",
        "Industry Recipe: DeepSeek",
        "DeepSeekMath 提出 GRPO，DeepSeek-R1 用 rule-based reward 和大规模 RL 展示 long-CoT emergence。",
        points=[
            "R1-Zero 强调纯 RL 对推理能力的激发。",
            "R1 用冷启动数据、RL、rejection sampling、distillation 改善可读性与通用性。",
            "经验启发是：数学代码等领域可以先不用神经 RM，而用严格 verifier 提供主奖励。",
        ],
        source="GRPO, R1",
    ),
    slide(
        "Qwen/GSPO 的重点是让 MoE 与长推理 RL 更稳定",
        "Industry Recipe: Qwen",
        "GSPO 把 clipping 从 token 提升到 sequence，是对 GRPO 在 MoE 场景中不稳定现象的直接回应。",
        points=[
            "sequence-level objective 更匹配 response-level reward。",
            "对 MoE 路由与 token-level probability noise 更鲁棒。",
            "这说明大厂 recipe 越来越重视 optimizer 与模型架构之间的耦合。",
        ],
        source="GSPO",
    ),
    slide(
        "ByteDance DAPO 的价值在于开源可复现的大规模 RLVR 工程闭环",
        "Industry Recipe: DAPO",
        "DAPO 与 verl 生态展示了从数据、采样、rollout、loss 到长度控制的一整套实践细节。",
        points=[
            "DAPO-Math-17K 表明少量高质量 prompt 配合多采样可驱动强推理 RL。",
            "四项修复都是围绕有效 token 梯度、熵、长度和难度分布。",
            "开源 recipe 的重要性在于让社区能复现实验，而不只看到 benchmark 分数。",
        ],
        source="DAPO",
    ),
    slide(
        "MiniMax-M1 的 CISPO 针对的是 off-policy 多轮复用下的梯度浪费",
        "Industry Recipe: MiniMax",
        "当一个 rollout batch 被更新很多次，PPO-style clip 会让越来越多 token 被截断；CISPO 保留了这些 token 的方向信息。",
        points=[
            "clipped IS weight 只修正采样分布差异，不把 log-prob 梯度直接置零。",
            "这对长链推理中的反思 token 尤其重要。",
            "它展示了推理 RL 的瓶颈不只是 reward，还包括训练数据复用率和 off-policy 修正。",
        ],
        source="MINIMAX",
    ),
    slide(
        "Kimi k1.5 把强化学习扩展到长上下文与长测试时计算",
        "Industry Recipe: Kimi",
        "Kimi k1.5 强调 long-CoT、long2short、数据筛选和 RL scaling，说明推理模型能力来自训练时与测试时计算的联合扩展。",
        points=[
            "长上下文让模型能承载更长思考与工具化中间状态。",
            "long2short 把长推理能力蒸馏到更短、更可用的回答中。",
            "这类 recipe 通常组合 SFT、RLVR、拒绝采样和蒸馏，而不是单一算法。",
        ],
        source="KIMI",
    ),
    slide(
        "Tulu 3 代表开放后训练 recipe 的系统化趋势",
        "Industry Recipe: Open Recipes",
        "开放模型越来越重视可复现数据混合、DPO/RLVR 配方、评测协议和训练脚本，而不只是发布 checkpoint。",
        points=[
            "Tulu 3、OpenRLHF、verl、slime 等项目把偏好优化和 RLVR 训练模块化。",
            "这降低了 DPO/GRPO/PPO 的工程门槛，也暴露了不同实现细节导致的结果差异。",
            "研究者应同时报告数据、采样策略、reward、KL、clip、batch token 和过滤规则。",
        ],
        source="TULU3, DAPO",
    ),
    slide(
        "评测必须区分训练 reward、真实偏好和任务正确性",
        "Evaluation",
        "一个算法在 reward 上升并不等于模型变好；RLHF/RLVR 中最常见的误判就是把代理指标当成目标。",
        table=[
            ["评测对象", "常见指标", "常见误区"],
            ["偏好质量", "Arena-Hard、AlpacaEval、人工 A/B", "judge 偏见、长度偏好、风格偏好"],
            ["推理正确性", "AIME、MATH、GSM8K、Codeforces、单测 pass rate", "污染、格式解析、采样预算不一致"],
            ["奖励模型", "RewardBench、held-out preference accuracy", "高准确率不代表可安全优化"],
            ["训练稳定性", "KL、entropy、clip fraction、response length、reward std", "只看最终分数忽略中途 collapse"],
        ],
        source="REWARDBENCH, R1, DAPO",
    ),
    slide(
        "reward hacking 通常先表现为格式、长度和裁判盲区",
        "Risk Control",
        "大模型不需要真正理解目标函数，只要找到能提高 reward 的可利用模式即可。",
        points=[
            "格式投机：输出满足 verifier parser 的模板，却没有真实推理。",
            "长度投机：用更长 CoT 增加蒙中概率，或用短答规避负 token 梯度。",
            "裁判投机：迎合 AI judge 的措辞、礼貌、冗长或安全模板。",
            "控制手段包括 KL、长度 reward、格式 parser、held-out verifier、多裁判交叉和人工抽检。",
        ],
        source="REWARDBENCH, DAPO, R1",
    ),
    slide(
        "关键超参数应该按可观测量闭环调节，而不是固定抄论文",
        "Hyperparameters",
        "β、ε、G、temperature、max length、batch token、update epochs 都通过 KL、entropy、clip fraction 和 reward variance 反馈。",
        points=[
            "β：observed KL 高于目标就增大，低于目标就减小。",
            "ε：clip fraction 长期过高说明步子太大或数据复用过多。",
            "G：组大小越大，advantage 越稳，但 rollout 成本线性增长。",
            "max length：不是越长越好，需要与 overlong penalty、推理预算和评测协议一致。",
        ],
    ),
    slide(
        "大规模 RLVR 的系统瓶颈是 rollout、验证器和 token 级吞吐",
        "Systems",
        "训练不再只是 GPU 上的反向传播；采样、过滤、verifier、分布式同步和 checkpoint 都会成为主路径。",
        points=[
            "rollout engine 需要高吞吐生成，常用 vLLM、SGLang 或自研服务。",
            "verifier 要低延迟且可复现，数学解析和代码沙箱尤其容易成为瓶颈。",
            "训练端要以 batch token 而非样本数估算负载，否则长 CoT 会打破吞吐计划。",
            "off-policy 数据复用提高效率，但必须用 IS、KL 或刷新 rollout 控制漂移。",
        ],
        source="DAPO, MINIMAX, TULU3",
    ),
    slide(
        "数据选择比 optimizer 更常决定 RL 是否学到推理而不是记忆模板",
        "Data Curation",
        "推理 RL 的 prompt 应有可验证答案、适中难度、多样解法和低污染风险。",
        points=[
            "太易样本产生全对组，太难样本产生全错组，二者都浪费 rollout。",
            "Dynamic Sampling 或 difficulty curriculum 可保持有效梯度密度。",
            "训练集要覆盖题型、语言、格式和长度，避免模型只学会某个 benchmark 的表面分布。",
            "偏好数据要保留 rejected answer，因为负例决定 margin 几何。",
        ],
        source="DAPO, R1, DPO",
    ),
    slide(
        "选择算法时先回答五个问题",
        "Decision Checklist",
        table=[
            ["问题", "若答案是是", "推荐路线"],
            ["能否程序验证正确性", "数学、代码、结构化任务", "GRPO/DAPO/GSPO/CISPO/RLOO"],
            ["已有大量偏好对", "离线人类或 AI preference", "DPO/IPO/KTO/SimPO/ORPO"],
            ["是否必须在线探索", "需要发现新解法或长 CoT", "PPO/RLOO/GRPO 系列"],
            ["是否能承担 critic", "reward dense 或 PRM 可靠", "PPO/VAPO"],
            ["是否只想稳定蒸馏", "筛选器强、算力有限", "RAFT/ReST/rejection SFT"],
        ],
    ),
    slide(
        "公式速查一：从 KL-RLHF 到 DPO 的核心等价",
        "Formula Sheet",
        points=[
            "KL-RLHF：最大化 reward，同时惩罚 policy 偏离 reference。",
            "BT-RM：用 winner-loser 奖励差拟合偏好概率。",
            "DPO：把最优策略的 log-ratio 作为隐式 reward，直接优化 preference logistic loss。",
            "IPO/KTO/SimPO/ORPO 都是在这个 log-ratio 或 log-prob 几何上改变 link、reference 或 margin。",
        ],
        formula="dpo",
        formula_caption="DPO 是偏好优化族的中心参照式。",
        source="DPO, IPO, KTO, ORPO, SIMPO",
    ),
    slide(
        "公式速查二：从 PPO 到 GRPO、DAPO、GSPO 的核心变化",
        "Formula Sheet",
        points=[
            "PPO：critic 给 token-level advantage，clip token-level ratio。",
            "GRPO：用同 prompt 多样本 reward 构造 group-relative advantage，去掉 value model。",
            "DAPO：改变 clip、采样、归一化和长度 shaping，提升有效训练信号。",
            "GSPO：把 ratio 与 clip 提升到 sequence level，匹配 response-level reward。",
        ],
        formula="gspo_ratio",
        formula_caption="GSPO 的 sequence ratio 是 GRPO 之后的重要结构性修改。",
        source="PPO, GRPO, DAPO, GSPO",
    ),
    slide(
        "研究时间线显示，领域从通用 RLHF 走向推理专用 RLVR",
        "Timeline",
        table=[
            ["阶段", "代表", "关键词"],
            ["2017-2022", "Human Preferences、InstructGPT、HH", "RM + PPO 标准化"],
            ["2023", "RRHF、RAFT、SLiC、DPO、ReST、ReMax、IPO", "不用在线 RL 的偏好优化爆发"],
            ["2024", "KTO、GRPO、RLOO、ORPO、SimPO、Online DPO", "critic-free 与 reference-free 两条线并进"],
            ["2025", "R1、REINFORCE++、DAPO、Dr.GRPO、VAPO、CISPO、GSPO", "RLVR、长 CoT、group/sequence objective 工程化"],
            ["2026", "持续演进", "多裁判、工具环境、agentic RL、跨模态 reward"],
        ],
    ),
    slide(
        "不要把算法名当作能力来源，能力来自目标、数据、采样和约束的闭环",
        "Synthesis",
        "DPO、PPO、GRPO、DAPO、GSPO 都只是策略改进算子的不同近似；真正要审查的是它优化了什么、在哪里采样、如何防止投机。",
        points=[
            "如果 reward 错，越强的 optimizer 越快放大错误。",
            "如果数据窄，online RL 也只会在窄分布内探索。",
            "如果 KL 与长度不可控，推理模型容易变成冗长且难评估的模型。",
            "专业后训练的核心能力是把公式、系统指标和人工审阅联成闭环。",
        ],
        callout="实践原则：先定义可信 reward，再选择最简单足够的 optimizer，最后用真实评测验证而不是用训练 reward 自证。",
    ),
    slide(
        "术语表把常见混淆拆开",
        "Glossary",
        table=[
            ["术语", "精确定义"],
            ["reference model", "KL 或 log-ratio 的参照策略，通常是 SFT 或前一阶段 policy"],
            ["old policy", "产生 rollout 的策略，用于 importance ratio"],
            ["reward model", "从偏好或规则估计标量 reward 的模型或函数"],
            ["verifier", "能判断输出是否满足规则的程序、模型或混合系统"],
            ["advantage", "相对 baseline 的收益增量，是 policy gradient 的权重"],
            ["clip fraction", "PPO/GRPO 中被剪切 token 或序列的比例，反映更新是否过激"],
        ],
    ),
    slide(
        "参考文献一：RLHF 与基础策略优化",
        "References",
        table=[
            ["ID", "Source"],
            ["TRPO", "Schulman et al., Trust Region Policy Optimization, 2015"],
            ["GAE", "Schulman et al., Generalized Advantage Estimation, 2015"],
            ["PPO", "Schulman et al., Proximal Policy Optimization Algorithms, 2017"],
            ["DRHF", "Christiano et al., Deep RL from Human Preferences, 2017"],
            ["IGPT", "Ouyang et al., InstructGPT, 2022"],
            ["HH/CAI", "Anthropic Helpful-Harmless RLHF and Constitutional AI, 2022"],
        ],
    ),
    slide(
        "参考文献二：偏好优化与 RLVR",
        "References",
        table=[
            ["ID", "Source"],
            ["DPO/IPO/KTO/ORPO/SimPO", "Direct and reference-free preference optimization family"],
            ["RLOO/ReMax/REINFORCE++", "critic-free policy gradient family"],
            ["GRPO/DeepSeek-R1", "group-relative RLVR and reasoning model training"],
            ["DAPO/Dr.GRPO/VAPO", "long-CoT RLVR stabilization and bias correction"],
            ["CISPO/GSPO", "off-policy IS clipping and sequence-level group optimization"],
            ["Tulu/Kimi/MiniMax/Qwen", "open and industrial post-training recipes"],
        ],
    ),
]


CSS = """
@font-face{font-family:Tsanger;src:url('fonts/TsangerJinKai02-W04.ttf') format('truetype');font-weight:400}
@font-face{font-family:Tsanger;src:url('fonts/TsangerJinKai02-W05.ttf') format('truetype');font-weight:700}
@font-face{font-family:JetBrainsMono;src:url('fonts/JetBrainsMono.woff2') format('woff2');font-weight:400}
@page{size:280mm 158mm;margin:0}
*{box-sizing:border-box}
html,body{margin:0;padding:0;background:#e9e5d8;color:#141413}
body{font-family:Georgia,Tsanger,"Noto Serif CJK SC","Source Han Serif SC",serif}
.latin{font-family:Arial,"Helvetica Neue",Helvetica,sans-serif;font-weight:inherit;letter-spacing:0}
.slide{position:relative;width:280mm;height:158mm;padding:13mm 16mm 11mm 18mm;background:#f5f4ed;break-after:page;overflow:hidden}
.slide::before{content:"";position:absolute;left:9mm;top:12mm;bottom:12mm;width:1.2mm;background:#1B365D}
.slide::after{content:"";position:absolute;left:0;right:0;bottom:0;height:2.2mm;background:#1B365D}
.eyebrow{font-family:JetBrainsMono,monospace;font-size:8.5pt;letter-spacing:0;color:#7B2D26;text-transform:uppercase;margin-bottom:3mm}
h1,h2{font-family:Georgia,Tsanger,"Noto Serif CJK SC","Source Han Serif SC",serif;font-weight:700;color:#1B365D;margin:0;line-height:1.08;letter-spacing:0}
h1{font-size:34pt;max-width:230mm}
h2{font-size:21.5pt;max-width:235mm}
.lead{font-size:13pt;line-height:1.45;margin:4mm 0 4.5mm;max-width:235mm;color:#2a2926}
.cover .lead{font-size:17pt;max-width:230mm;margin-top:7mm}
.bullets{margin:2mm 0 0 0;padding:0;list-style:none;display:grid;gap:2.4mm;max-width:236mm}
.bullets li{position:relative;font-size:10.9pt;line-height:1.44;padding-left:5.5mm}
.bullets li::before{content:"";position:absolute;left:0;top:.64em;width:2.1mm;height:2.1mm;background:#7B2D26}
.cols{display:grid;grid-template-columns:repeat(3,1fr);gap:5mm;margin-top:5mm}
.col{border-top:1.2pt solid #1B365D;padding-top:3mm}
.col h3{font-size:13pt;margin:0 0 2mm;color:#1B365D}
.col p,.col li{font-size:10.2pt;line-height:1.42}
.col ul{list-style:none;margin:0;padding:0;display:grid;gap:2mm}
.formula-card{margin:4mm 0 2.8mm;padding:4mm 5mm;background:#fffdf6;border:1.2pt solid rgba(27,54,93,.38);border-radius:2mm;display:flex;align-items:center;justify-content:center;min-height:20mm;max-height:37mm}
.formula-card img{max-width:100%;max-height:30mm;object-fit:contain}
.caption{font-size:9.2pt;line-height:1.35;color:#5c5547;margin:0 0 3mm;font-family:JetBrainsMono,Tsanger,monospace}
table{width:100%;border-collapse:collapse;margin-top:4mm;font-size:8.9pt;line-height:1.32}
th,td{border-bottom:1px solid rgba(27,54,93,.28);padding:2.0mm 2.1mm;text-align:left;vertical-align:top}
th{color:#1B365D;font-weight:700;background:#ece7d8}
tr:first-child td{color:#1B365D;font-weight:700;background:#ece7d8}
.callout{position:absolute;left:18mm;right:16mm;bottom:10mm;border-left:3mm solid #7B2D26;background:#fffdf6;padding:3mm 4mm;font-size:11.4pt;line-height:1.42;color:#1d1c19}
.source{position:absolute;right:16mm;bottom:5.5mm;font-family:JetBrainsMono,monospace;font-size:7pt;color:#6d685f;max-width:160mm;text-align:right}
.page{position:absolute;left:18mm;bottom:5.5mm;font-family:JetBrainsMono,monospace;font-size:7pt;color:#6d685f}
.cover{padding-top:20mm}
.cover::before{width:1.8mm}
.cover h1{font-size:42pt}
.cover .bullets{margin-top:7mm;max-width:215mm}
.cover .bullets li{font-size:12.4pt}
"""


def render_table(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    body = []
    for row in rows:
        body.append("<tr>" + "".join(f"<td>{fmt(cell)}</td>" for cell in row) + "</tr>")
    return "<table>" + "".join(body) + "</table>"


def render_columns(columns: list[tuple[str, list[str]]]) -> str:
    if not columns:
        return ""
    out = ["<div class='cols'>"]
    for title, items in columns:
        out.append("<div class='col'>")
        out.append(f"<h3>{fmt(title)}</h3>")
        out.append("<ul>")
        for item in items:
            out.append(f"<li>{fmt(item)}</li>")
        out.append("</ul></div>")
    out.append("</div>")
    return "".join(out)


def render_slide(data: dict[str, Any], idx: int) -> str:
    kind = esc(data.get("kind", "standard"))
    title_tag = "h1" if kind == "cover" else "h2"
    out = [f"<section class='slide {kind}'>"]
    out.append(f"<div class='eyebrow'>{esc(data['eyebrow'])}</div>")
    out.append(f"<{title_tag}>{fmt(data['title'])}</{title_tag}>")
    if data.get("lead"):
        out.append(f"<p class='lead'>{fmt(data['lead'])}</p>")
    if data.get("formula"):
        key = data["formula"]
        out.append(f"<div class='formula-card'><img src='assets/formulas/{esc(key)}.svg' alt='LaTeX formula {esc(key)}'></div>")
        if data.get("formula_caption"):
            out.append(f"<p class='caption'>{fmt(data['formula_caption'])}</p>")
    if data.get("columns"):
        out.append(render_columns(data["columns"]))
    if data.get("points"):
        out.append("<ul class='bullets'>")
        for item in data["points"]:
            out.append(f"<li>{fmt(item)}</li>")
        out.append("</ul>")
    if data.get("table"):
        out.append(render_table(data["table"]))
    if data.get("callout"):
        out.append(f"<div class='callout'>{fmt(data['callout'])}</div>")
    if data.get("source"):
        out.append(f"<div class='source'>Sources: {esc(data['source'])}</div>")
    out.append(f"<div class='page'>{idx:02d} / {len(SLIDES):02d}</div>")
    out.append("</section>")
    return "\n".join(out)


def write_html() -> None:
    html_doc = f"""<!doctype html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{esc(TITLE)}</title>
<meta name="author" content="Codex">
<meta name="date" content="{TODAY}">
<style>{CSS}</style>
</head>
<body>
{''.join(render_slide(s, i + 1) for i, s in enumerate(SLIDES))}
</body>
</html>
"""
    (ROOT / "index.html").write_text(html_doc, encoding="utf-8")


def write_sources() -> None:
    rows = ["# Sources\n", "本课件截至 2026-05-16 整理。公式与算法解读以论文原文和公开技术报告为主要来源。\n"]
    rows.append("| ID | Source | URL | Used for |\n|---|---|---|---|\n")
    for src in SOURCES:
        rows.append(f"| {src['id']} | {src['title']} | {src['url']} | {src['used']} |\n")
    (ROOT / "sources.md").write_text("".join(rows), encoding="utf-8")

    tsv = ["id\ttitle\turl\tused_for\n"]
    for src in SOURCES:
        tsv.append(f"{src['id']}\t{src['title']}\t{src['url']}\t{src['used']}\n")
    (ROOT / "data" / "source-map.tsv").write_text("".join(tsv), encoding="utf-8")


def write_readme() -> None:
    readme = f"""# {TITLE}

{SUBTITLE}

This directory contains an independently designed Chinese slide deck on reinforcement learning algorithms for large language model post-training. It covers RLHF, RLAIF, RLVR, PPO, RLOO, ReMax, REINFORCE++, DPO, IPO, KTO, ORPO, SimPO, GRPO, Dr.GRPO, DAPO, GSPO, CISPO, VAPO, RAFT, ReST, RRHF, SLiC-HF and representative industrial recipes.

## Artifacts

- `index.html` — source slide deck for browser and PDF rendering.
- `{SLUG}.pdf` — printable PDF deck.
- `{SLUG}.pptx` — PowerPoint export with rendered LaTeX formula images.
- `sources.md` — source map and references.
- `assets/formulas/` — LaTeX-generated SVG and PNG formula assets.

## Build

```bash
python3 build_deck.py
python3 -c "from weasyprint import HTML; HTML('index.html').write_pdf('{SLUG}.pdf')"
```

The formula assets are compiled with `latex` and `dvisvgm`; PPTX export uses `python-pptx`.
"""
    (ROOT / "README.md").write_text(readme, encoding="utf-8")


def text_lines_for_ppt(slide_data: dict[str, Any]) -> list[str]:
    lines: list[str] = []
    if slide_data.get("lead"):
        lines.append(slide_data["lead"])
    for title, items in slide_data.get("columns", []):
        lines.append(f"{title}: " + "；".join(items))
    for item in slide_data.get("points", []):
        lines.append(item)
    if slide_data.get("table"):
        for row in slide_data["table"][1:]:
            lines.append(" | ".join(row))
    if slide_data.get("callout"):
        lines.append(slide_data["callout"])
    return lines[:9]


def write_pptx() -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_AUTO_SIZE
        from pptx.util import Inches, Pt
    except Exception as exc:
        print(f"[pptx skipped] python-pptx not available: {exc}")
        return

    prs = Presentation()
    prs.slide_width = Inches(11.02)
    prs.slide_height = Inches(6.22)
    blank = prs.slide_layouts[6]
    for idx, data in enumerate(SLIDES, start=1):
        sl = prs.slides.add_slide(blank)
        bg = sl.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(245, 244, 237)

        bar = sl.shapes.add_shape(1, Inches(0.35), Inches(0.42), Inches(0.04), Inches(5.35))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(27, 54, 93)
        bar.line.fill.background()

        eyebrow = sl.shapes.add_textbox(Inches(0.65), Inches(0.36), Inches(9.7), Inches(0.25)).text_frame
        eyebrow.text = data["eyebrow"]
        eyebrow.paragraphs[0].font.name = "JetBrains Mono"
        eyebrow.paragraphs[0].font.size = Pt(8)
        eyebrow.paragraphs[0].font.color.rgb = RGBColor(123, 45, 38)

        title = sl.shapes.add_textbox(Inches(0.65), Inches(0.72), Inches(9.8), Inches(0.9)).text_frame
        title.text = data["title"]
        title.paragraphs[0].font.name = "Georgia"
        title.paragraphs[0].font.bold = True
        title.paragraphs[0].font.size = Pt(27 if data.get("kind") == "cover" else 19)
        title.paragraphs[0].font.color.rgb = RGBColor(27, 54, 93)

        y_top = 1.55 if not data.get("formula") else 1.42
        if data.get("formula"):
            img = ROOT / "assets" / "formulas" / f"{data['formula']}.png"
            if img.exists():
                sl.shapes.add_picture(str(img), Inches(0.75), Inches(y_top), width=Inches(9.6))
                y_top += 1.15
        box = sl.shapes.add_textbox(Inches(0.75), Inches(y_top), Inches(9.75), Inches(4.15 - max(0, y_top - 1.55))).text_frame
        box.word_wrap = True
        box.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        box.clear()
        lines = text_lines_for_ppt(data)
        if not lines:
            lines = [data.get("formula_caption", "")]
        for i, line in enumerate(lines):
            p = box.paragraphs[0] if i == 0 else box.add_paragraph()
            p.text = line
            p.level = 0
            p.font.name = "Georgia"
            p.font.size = Pt(11 if len(lines) > 6 else 12.5)
            p.font.color.rgb = RGBColor(20, 20, 19)

        foot = sl.shapes.add_textbox(Inches(0.65), Inches(5.82), Inches(2.0), Inches(0.2)).text_frame
        foot.text = f"{idx:02d} / {len(SLIDES):02d}"
        foot.paragraphs[0].font.name = "JetBrains Mono"
        foot.paragraphs[0].font.size = Pt(7)
        foot.paragraphs[0].font.color.rgb = RGBColor(100, 96, 88)

    prs.save(ROOT / f"{SLUG}.pptx")


def write_pdf() -> None:
    from weasyprint import HTML

    HTML(str(ROOT / "index.html")).write_pdf(str(ROOT / f"{SLUG}.pdf"))


def main() -> None:
    (ROOT / "assets" / "formulas").mkdir(parents=True, exist_ok=True)
    (ROOT / "data").mkdir(parents=True, exist_ok=True)
    render_formulas()
    write_html()
    write_sources()
    write_readme()
    write_pdf()
    write_pptx()


if __name__ == "__main__":
    main()
